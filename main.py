import os
import aiohttp
from dotenv import load_dotenv
from pyrogram import Client, filters
from pyrogram.types import Message, InlineKeyboardMarkup, InlineKeyboardButton
from openai import OpenAI
from bs4 import BeautifulSoup
from docx import Document
from pptx import Presentation
import openai
import base64
import mimetypes
from pyrogram.enums import ParseMode
import datetime
from sqlalchemy import create_engine, Column, Integer, String, Text, DateTime, ForeignKey, JSON
from sqlalchemy.orm import sessionmaker, relationship, declarative_base, scoped_session
from sqlalchemy.future import select # Используем select из sqlalchemy.future для совместимости
import asyncio # Нужен для ожидания
import aiofiles

single_document_filter = filters.document & ~filters.media_group
single_photo_filter = filters.photo & ~filters.media_group

# Вверху файла:
user_states = {}

# Словарь для буферизации медиа-групп: {media_group_id: {'messages': [msg1, msg2,...], 'timer_task': task_handle}}
media_group_buffers = {}
# Блокировка для безопасного доступа к буферу из разных обработчиков (на всякий случай)
buffer_lock = asyncio.Lock()
# Задержка в секундах перед обработкой группы (чтобы успели прийти все сообщения)
MEDIA_GROUP_DELAY = 2.0
# Telegram max length per message
MAX_LENGTH = 4096

# Загрузка переменных из .env
load_dotenv()
API_ID = int(os.getenv("API_ID"))
API_HASH = os.getenv("API_HASH")
BOT_TOKEN = os.getenv("BOT_TOKEN")
OPENAI_API_KEY = os.getenv("OPENAI_API_KEY")
GEMINI_API_KEY = os.getenv("GOOGLE_API")
DEEPSEEK = os.getenv("DEEPSEEK_API")
GROQ_API = os.getenv("GROQ_API")
GROK_API = os.getenv("GROK_API")
GLM_API = os.getenv("GLM_API")

# Файлы и переменные
MODEL_CATEGORIES = {
    "openai": {
        "title": "🧠 Модели OpenAI",
        "models": {
            "gpt-4o-mini": "GPT-4o mini ",
            "gpt-4o": "GPT-4o",
            "gpt-4.1": "GPT-4.1",
            "gpt-4.1-mini": "GPT-4.1 Mini",
            "gpt-4.1-nano": "GPT-4.1 Nano",
            "o4-mini": "o4 Mini",
            "o3": "o3",
            "o3-pro": "o3-pro"
        }
    },
    "gemini": {
        "title": "🚀 Модели Gemini",
        "models": {
            "gemini-2.5-pro": "Gemini 2.5 pro",
            "gemini-2.5-flash": "Gemini 2.5 Flash  🚀",
            "gemini-2.0-flash": "Gemini 2.0 flash 🌟",
            "gemini-2.0-flash-lite": "Gemini 2.0 flash Lite ⚡"
        }
    },
    "deepseek": {
        "title": "🧪 Модели DeepSeek",
        "models": {
            "deepseek-chat": "DeepSeek Chat 🤖",
            "deepseek-reasoner": "DeepSeek Reasoner 🧠",
            "deepseek-r1-distill-llama-70b": "DeepSeek R1 Distill 🧪🦙"
        }
    },
    "llama": {
        "title": "🦙 Модели LLaMA & Mistral",
        "models": {
            "mistral-saba-24b": "Mistral Saba 24B 🌬️",
            "meta-llama/llama-4-maverick-17b-128e-instruct": "LLaMA 4 maverick 🦙🔥",
            "compound-beta": "Compound beta",
            "compound-beta-mini": "Compound beta mini"
        }
    },
    "Grok": {
        "title": "🏆 Grok",
        "models": {
            "grok-4-0709": "Grok 4"
        }
    },
    "GLM": {
        "title": "GLM Z",
        "models": {
            "GLM-4.5": "GLM 4.5",
            "GLM-4.5-X": "GLM 4.5 X",
            "GLM-4.5-Air": "GLM 4.5 Air"
        }
    }
}


GOOGLE_MODELS = ["gemini-2.5-pro", "gemini-2.5-flash", "gemini-2.0-flash", "gemini-2.0-flash-lite"]
DEEPSEEK_MODELS = ["deepseek-chat", "deepseek-reasoner"]
GROQ_MODELS = ["mistral-saba-24b", "meta-llama/llama-4-maverick-17b-128e-instruct", "deepseek-r1-distill-llama-70b", "compound-beta", "compound-beta-mini"]
GROK_MODELS = ["grok-4-0709"]
GLM_MODELS = ["GLM-4.5-Air", "GLM-4.5-X", "GLM-4.5"]

# Настройка клиентов
client = Client("GPTBot", api_id=API_ID, api_hash=API_HASH, bot_token=BOT_TOKEN)
openai.api_key = OPENAI_API_KEY
client_ai = OpenAI()
client_deepseek = OpenAI(
    base_url="https://api.deepseek.com",
    api_key=DEEPSEEK
    )

client_google = OpenAI(
    api_key=GEMINI_API_KEY,
    base_url="https://generativelanguage.googleapis.com/v1beta/openai/"
    )

client_groq = OpenAI(
    base_url="https://api.groq.com/openai/v1",
    api_key=GROQ_API
)

client_grok = OpenAI(
    base_url = "https://api.x.ai/v1",
    api_key=GROK_API
)

client_glm = OpenAI(
    base_url = "https://api.z.ai/api/paas/v4",
    api_key=GLM_API
)
state = {"client_now":client_ai}

# Функция для обработки собранной медиа-группы
async def process_media_group(media_group_id: str, chat_id: str, client_instance):
    async with buffer_lock:
        if media_group_id not in media_group_buffers:
            return # Группа уже обработана или произошла ошибка

        grouped_messages = media_group_buffers[media_group_id]['messages']
        # Убираем группу из буфера сразу, чтобы избежать повторной обработки
        del media_group_buffers[media_group_id]

    if not grouped_messages:
        return

    # Определяем тип медиа (фото или документы) и ищем подпись
    first_message = grouped_messages[0]
    caption = first_message.caption or "" # Подпись обычно у первого сообщения
    combined_content = None
    message_to_reply = first_message # Будем отвечать на первое сообщение группы

    # --- Обработка Группы Изображений ---
    if first_message.photo:
        image_contents = []
        if caption: # Текст добавляем первым элементом
             image_contents.append({"type": "text", "text": caption})
        else: # Добавляем текст по умолчанию, если нет подписи
             image_contents.append({"type": "text", "text": ""})


        processing_message = await client_instance.send_message(chat_id, f"⏳ Обрабатываю {len(grouped_messages)} изображений...")

        temp_files = []
        try:
            for msg in grouped_messages:
                if msg.photo:
                    file_path = await msg.download()
                    temp_files.append(file_path)
                    encoded_url = encode_image_as_base64(file_path) # Ваша функция кодирования
                    image_contents.append({"type": "image_url", "image_url": {"url": encoded_url}})

            combined_content = image_contents # Финальный контент - список словарей

        except Exception as e:
            await processing_message.edit_text(f"❌ Ошибка при загрузке/кодировании изображений: {e}")
            return # Прерываем обработку
        finally:
            # Удаляем временные файлы
            for f_path in temp_files:
                if os.path.exists(f_path):
                    os.remove(f_path)
            # Удаляем сообщение "Обрабатываю..."
            await processing_message.delete()


    # --- Обработка Группы Документов ---
    elif first_message.document:
        all_texts = []
        doc_names = []
        processing_message = await client_instance.send_message(chat_id, f"⏳ Обрабатываю {len(grouped_messages)} документов...")
        temp_files = []
        try:
            for msg in grouped_messages:
                if msg.document:
                    file_path = await msg.download()
                    temp_files.append(file_path)
                    doc_names.append(os.path.basename(file_path))
                    ext = os.path.splitext(file_path)[1].lower()
                    doc_text = None

                    # Логика извлечения текста (аналогично handle_file)
                    if ext == ".txt":
                        async with aiofiles.open(file_path, "r", encoding="utf-8") as f:
                            doc_text = await f.read()
                    elif ext == ".docx":
                        doc = Document(file_path)
                        doc_text = "\n".join([p.text for p in doc.paragraphs])
                    elif ext == ".pptx":
                        prs = Presentation(file_path)
                        lines = [shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text")]
                        doc_text = "\n".join(lines)
                    elif ext == ".fb2":
                        async with aiofiles.open(file_path, "r", encoding="utf-8") as f:
                             content = await f.read()
                             soup = BeautifulSoup(content, "lxml")
                             doc_text = soup.get_text(separator="\n", strip=True)
                    # Добавьте другие поддерживаемые форматы, если нужно

                    if doc_text:
                        # Ограничение длины текста одного документа
                        if len(doc_text) > 20000:
                             doc_text = doc_text[:20000] + "\n... (текст документа обрезан)"
                        all_texts.append(f"\n\n--- Текст из файла {os.path.basename(file_path)} ---\n{doc_text}")
                    else:
                         all_texts.append(f"\n\n--- Не удалось извлечь текст из файла {os.path.basename(file_path)} ({ext}) ---")


            # Собираем финальный текст
            combined_text = (caption + "\n" if caption else "")
            combined_text += f"[Обработаны файлы: {', '.join(doc_names)}]"
            combined_text += "".join(all_texts)

            # Ограничение общей длины текста
            if len(combined_text) > 40000:
                 combined_text = combined_text[:40000] + "\n... (общий текст файлов обрезан)"

            combined_content = combined_text # Финальный контент - строка

        except Exception as e:
            await processing_message.edit_text(f"❌ Ошибка при обработке документов: {e}")
            return
        finally:
             for f_path in temp_files:
                 if os.path.exists(f_path):
                     os.remove(f_path)
             await processing_message.delete()

    # --- Вызов основной логики обработки ---
    if combined_content:
        # Вызываем вашу основную функцию process_message, передавая первое сообщение группы
        # (для ответа на него) и собранный контент
        await process_message(message_to_reply, combined_content)


# Обработчик для сообщений, входящих в медиа-группу
@client.on_message(filters.media_group)

async def media_group_handler(c: Client, message: Message):
    media_group_id = message.media_group_id
    chat_id = str(message.chat.id)

    async with buffer_lock:
        # Если это первое сообщение группы
        if media_group_id not in media_group_buffers:
            media_group_buffers[media_group_id] = {'messages': [message], 'timer_task': None}
            # Запускаем таймер на обработку
            task = asyncio.create_task(
                schedule_group_processing(media_group_id, chat_id, c)
            )
            media_group_buffers[media_group_id]['timer_task'] = task
        else:
            # Добавляем сообщение в существующую группу
            media_group_buffers[media_group_id]['messages'].append(message)
            # Перезапускаем таймер (отменяем старый, создаем новый)
            if media_group_buffers[media_group_id]['timer_task']:
                media_group_buffers[media_group_id]['timer_task'].cancel()
            task = asyncio.create_task(
                schedule_group_processing(media_group_id, chat_id, c)
            )
            media_group_buffers[media_group_id]['timer_task'] = task

# Вспомогательная функция для отложенного запуска
async def schedule_group_processing(media_group_id: str, chat_id: str, client_instance):
    await asyncio.sleep(MEDIA_GROUP_DELAY)
    # Проверяем еще раз под блокировкой, существует ли группа (вдруг ее уже обработали или отменили)
    async with buffer_lock:
        if media_group_id in media_group_buffers:
             # Запускаем основную обработку группы
             asyncio.create_task(process_media_group(media_group_id, chat_id, client_instance))

# Настройка базы данных
DATABASE_URL = "sqlite:///gpt_bot_data.db" # Файл базы данных будет создан в той же папке
engine = create_engine(DATABASE_URL, echo=False) # echo=True для отладки SQL запросов
Base = declarative_base()
SessionLocal = scoped_session(sessionmaker(autocommit=False, autoflush=False, bind=engine))

# Определяем модели таблиц
class Chat(Base):
    __tablename__ = "chats"
    chat_id = Column(String, primary_key=True, index=True) # Используем String для chat_id
    model_name = Column(String, default="gpt-4o-mini")
    system_prompt = Column(Text, default="")
    # Связь с сообщениями (для удобства, но не обязательно для основного функционала)
    messages = relationship("Message", back_populates="chat", cascade="all, delete-orphan")

class Message(Base):
    __tablename__ = "messages"
    id = Column(Integer, primary_key=True, index=True)
    chat_id = Column(String, ForeignKey("chats.chat_id"), index=True)
    role = Column(String) # 'system', 'user', 'assistant'
    content = Column(JSON) # Используем JSON для хранения как текста, так и сложных структур (для картинок)
    timestamp = Column(DateTime, default=datetime.datetime.utcnow)
    # Связь с чатом
    chat = relationship("Chat", back_populates="messages")

# Создаем таблицы, если их нет
Base.metadata.create_all(bind=engine)

# Функция для получения сессии базы данных (управляет сессиями)
def get_db():
    db = SessionLocal()
    try:
        yield db
    finally:
        db.close()



def encode_image_as_base64(path: str):
    with open(path, "rb") as f:
        encoded = base64.b64encode(f.read()).decode("utf-8")
    mime_type, _ = mimetypes.guess_type(path)
    return f"data:{mime_type};base64,{encoded}"




# Команда /start
@client.on_message(filters.command("start"))
async def start(_, message: Message):
    chat_id = str(message.chat.id)
    db = next(get_db()) # Получаем сессию
    try:
        # Проверяем, существует ли чат
        existing_chat = db.execute(select(Chat).filter(Chat.chat_id == chat_id)).scalar_one_or_none()

        if existing_chat is None:
            # Создаем новый чат
            new_chat = Chat(
                chat_id=chat_id,
                model_name="gpt-4.1-mini", # Модель по умолчанию
                system_prompt=""          # Пустой системный промпт по умолчанию
            )
            db.add(new_chat)
            db.commit() # Сохраняем изменения в БД
            await message.reply_text("✅ Чат зарегистрирован для общения с GPT.")
        else:
            await message.reply_text("⚠️ Этот чат уже зарегистрирован.")
    finally:
        db.close() # Закрываем сессию

from pyrogram.types import BotCommand

@client.on_message(filters.command("update_commands"))
async def update_commands(_, message: Message):
    commands = [
        BotCommand("start", "Регистрация чата для общения с GPT"),
        BotCommand("forget", "Сбросить память (контекст)"),
        BotCommand("context", "Установить новый системный контекст"),
        BotCommand("model", "Выбрать модель GPT"),
        BotCommand("gen", "Сгенерировать изображение DALL·E"),
        BotCommand("info", "Показать текущую модель и системный промпт"),
        BotCommand("help", "Показать справочное меню"),
        BotCommand("reset_context", "Удалить системный промпт")
    ]
    await client.set_bot_commands(commands)
    await message.reply_text("✅ Команды успешно обновлены!")


# Команда /forget
@client.on_message(filters.command("forget"))
async def forget(_, message: Message):
    chat_id = str(message.chat.id)
    db = next(get_db())
    try:
        chat = db.execute(select(Chat).filter(Chat.chat_id == chat_id)).scalar_one_or_none()
        if chat:
            # Удаляем все сообщения для этого чата
            db.query(Message).filter(Message.chat_id == chat_id).delete()
            # Опционально: сбрасываем системный промпт, если нужно
            # chat.system_prompt = ""
            db.commit()
            await message.reply_text("🧹 Контекст (история сообщений) очищен.")
        else:
            await message.reply_text("❗ Чат не зарегистрирован. Напиши /start.")
    finally:
        db.close()

@client.on_message(filters.command("context"))
async def ask_context(_, message: Message):
    user_states[message.chat.id] = "awaiting_context"

    cancel_button = InlineKeyboardMarkup([
        [InlineKeyboardButton("❌ Отменить", callback_data="cancel_context")]
    ])

    await message.reply_text("🧠 Введи новый системный контекст для этого чата.", reply_markup=cancel_button)

@client.on_message(filters.command("reset_context"))
async def reset_context(_, message: Message):
    chat_id = str(message.chat.id)
    db = next(get_db())

    try:
        chat = db.execute(select(Chat).filter(Chat.chat_id == chat_id)).scalar_one_or_none()

        if chat is None:
            await message.reply_text("❗ Чат не зарегистрирован. Напиши /start.")
            return

        # Сброс только system_prompt
        chat.system_prompt = ""
        db.commit()

        await message.reply_text("🧠 Системный контекст очищен. Теперь бот будет отвечать без специального поведения.")
    finally:
        db.close()

# Команда /model с кнопками
@client.on_message(filters.command("model"))
async def choose_model(_, message: Message):
    buttons = [
        [InlineKeyboardButton(cat_data["title"], callback_data=f"cat:{cat_key}")]
        for cat_key, cat_data in MODEL_CATEGORIES.items()
    ]
    await message.reply(
        "📂 Выбери категорию моделей:",
        reply_markup=InlineKeyboardMarkup(buttons)
    )

@client.on_callback_query(filters.regex(r"cat:(.+)"))
async def show_models_in_category(_, query):
    cat_key = query.data.split(":")[1]
    cat_data = MODEL_CATEGORIES.get(cat_key)

    if not cat_data:
        await query.answer("❌ Категория не найдена", show_alert=True)
        return

    model_buttons = [
        [InlineKeyboardButton(text=model_title, callback_data=f"model:{model_code}")]
        for model_code, model_title in cat_data["models"].items()
    ]
    model_buttons.append([InlineKeyboardButton("🔙 Назад", callback_data="back_to_categories")])

    await query.edit_message_text(
        f"{cat_data['title']}\n\nВыбери модель:",
        reply_markup=InlineKeyboardMarkup(model_buttons)
    )

@client.on_callback_query(filters.regex("back_to_categories"))
async def back_to_categories(_, query):
    buttons = [
        [InlineKeyboardButton(cat_data["title"], callback_data=f"cat:{cat_key}")]
        for cat_key, cat_data in MODEL_CATEGORIES.items()
    ]
    await query.edit_message_text(
        "📂 Выбери категорию моделей:",
        reply_markup=InlineKeyboardMarkup(buttons)
    )


@client.on_callback_query(filters.regex(r"model:"))
async def model_callback(_, query):
    model_name = query.data.split(":")[1]
    chat_id = str(query.message.chat.id)

    # Находим "красивое" имя модели по словарю
    pretty_name = next(
        (title for cat in MODEL_CATEGORIES.values() for code, title in cat["models"].items() if code == model_name),
        model_name  # fallback, если не найдём
    )

    db = next(get_db())
    try:
        chat = db.execute(select(Chat).filter(Chat.chat_id == chat_id)).scalar_one_or_none()
        if chat is None:
            await query.answer("Сначала зарегистрируй чат: /start", show_alert=True)
            return

        if model_name in GOOGLE_MODELS:
            state["client_now"] = client_google
        elif model_name in DEEPSEEK_MODELS:
            state["client_now"] = client_deepseek
        elif model_name in GROQ_MODELS:
            state["client_now"] = client_groq
        elif model_name in GROK_MODELS:
            state["client_now"] = client_grok
        elif model_name in GLM_MODELS:
            state["client_now"] = client_glm
        else:
            state["client_now"] = client_ai

        chat.model_name = model_name
        db.commit()

        await query.edit_message_text(f"✅ Модель установлена: <b>{pretty_name}</b>", parse_mode=ParseMode.HTML)

    finally:
        db.close()


# Автоответ на текст
# Общая функция для обработки сообщений (текст, файл, картинка)
async def process_message(message: Message, user_content: any):
    chat_id = str(message.chat.id)
    db = next(get_db())
    try:
        chat = db.execute(select(Chat).filter(Chat.chat_id == chat_id)).scalar_one_or_none()
        if chat is None:
            # Если чат не зарегистрирован, можно либо игнорировать, либо ответить
            # await message.reply_text("❗ Чат не зарегистрирован. Напиши /start.")
            return

        

        # 2. Формируем историю для API
        history_for_api = []
        # Добавляем системный промпт, если он есть
        if chat.system_prompt:
            history_for_api.append({"role": "system", "content": chat.system_prompt})

        # Получаем последние N сообщений из БД (например, 20)
        # Важно ограничить историю, чтобы не превышать лимиты токенов API
        # Учитываем, что content хранится как JSON
        db_messages = db.execute(
            select(Message.role, Message.content)
            .filter(Message.chat_id == chat_id)
            .order_by(Message.timestamp.desc()) # Сначала новые
            .limit(20) # Ограничиваем количество сообщений
        ).fetchall()

        # Добавляем сообщения в историю в правильном порядке (старые -> новые)
        for role, content in reversed(db_messages): # Переворачиваем для правильного порядка
             history_for_api.append({"role": role, "content": content})

        history_for_api.append({"role": "user", "content": user_content})
        # 3. Отправляем запрос к OpenAI
        try:
            client_now = state["client_now"]
            resp = client_now.chat.completions.create(
                model=chat.model_name, # Берем модель из настроек чата
                messages=history_for_api
            )
            reply_content = resp.choices[0].message.content
            if not reply_content:
                raise ValueError("Эта модель не может ответить\nпопробуйте сменить модель /model \nили очистить историю /forget")
            # 4. Сохраняем ответ ассистента в БД
            # 1. Сохраняем сообщение пользователя в БД
            
            db.add(Message(chat_id=chat_id, role="user", content=user_content))
            db.add(Message(chat_id=chat_id, role="assistant", content=reply_content))
            db.commit()

            # Разбиваем длинный ответ на части
            if len(reply_content) <= MAX_LENGTH:
                await message.reply_text(reply_content)
            else:
                parts = [reply_content[i:i+MAX_LENGTH] for i in range(0, len(reply_content), MAX_LENGTH)]
                for i, part in enumerate(parts):
                    await message.reply_text(part)

        except Exception as e:
            await message.reply_text(f"❌ Ошибка OpenAI: {e}")
            # Можно добавить откат последнего сообщения пользователя, если API упал
            # db.rollback()

    finally:
        db.close()

# Адаптируем существующие хендлеры, чтобы они вызывали process_message

@client.on_message(filters.text & ~filters.command(["start", "forget", "context", "model", "gen", "info", "help"]))
async def chat_handler(_, message: Message):
    chat_id = message.chat.id
    state = user_states.get(chat_id)

    if state == "awaiting_context":
        db = next(get_db())
        try:
            chat = db.execute(select(Chat).filter(Chat.chat_id == str(chat_id))).scalar_one_or_none()
            if chat:
                chat.system_prompt = message.text
                db.commit()
                await message.reply_text("✅ Новый системный контекст установлен.")
            else:
                await message.reply_text("❗ Сначала напиши /start.")
        finally:
            db.close()
        user_states.pop(chat_id)
        return

    elif state == "awaiting_prompt":
        prompt = message.text
        imamess = await message.reply_text("🎨 Генерирую изображение...")

        try:
            image = client_ai.images.generate(
                model="dall-e-3",
                prompt=prompt,
                n=1,
                size="1024x1024",
                quality="standard"
            )
            url = image.data[0].url

            await imamess.edit_text(
                f"🖼 <b>Готово!</b>\n"
                f"prompt: <code>{image.data[0].revised_prompt}</code>\n\n"
                f"<a href='{url}'>🔗 Открыть картинку</a>",
                parse_mode=ParseMode.HTML,
                disable_web_page_preview=False  # Показываем превью картинки
            )

        except Exception as e:
            await imamess.edit_text(f"❌ Ошибка генерации: {e}")

        user_states.pop(chat_id)
        return

    # Обычный текст — к GPT
    await process_message(message, message.text)



@client.on_message(single_document_filter)
async def handle_file(_, message: Message):
    chat_id = str(message.chat.id)
    # Проверка регистрации чата происходит внутри process_message
    # Но скачивание и обработка файла остаются здесь

    file_path = await message.download()
    ext = os.path.splitext(file_path)[1].lower()
    text = None

    try:
        if ext == ".txt":
            # Используем aiofiles для асинхронного чтения
            async with aiofiles.open(file_path, "r", encoding="utf-8") as f:
                text = await f.read()
        elif ext == ".docx":
            doc = Document(file_path)
            text = "\n".join([p.text for p in doc.paragraphs])
        elif ext == ".pptx":
            prs = Presentation(file_path)
            text_lines = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text_lines.append(shape.text)
            text = "\n".join(text_lines)
        elif ext == ".fb2":
             async with aiofiles.open(file_path, "r", encoding="utf-8") as f:
                content = await f.read()
                soup = BeautifulSoup(content, "lxml") # или 'html.parser'
                text = soup.get_text(separator="\n", strip=True)
        else:
            await message.reply_text("❌ Формат файла не поддерживается.")
            return

        if text:
            # Ограничим длину текста, если нужно
            if len(text) > 40000: # Примерное ограничение
                text = text[:40000] + "\n... (текст файла обрезан)"

            user_content = f"{message.caption or ''}\n[Содержимое файла {os.path.basename(file_path)}]:\n{text}"
            await process_message(message, user_content) # Передаем обработанный текст

    except Exception as e:
        await message.reply_text(f"❌ Ошибка обработки файла: {e}")
    finally:
        if os.path.exists(file_path):
            os.remove(file_path)


@client.on_message(single_photo_filter)
async def handle_base64_image(_, message: Message):
    chat_id = str(message.chat.id)
    # Проверка регистрации происходит внутри process_message

    file_path = None
    try:
        file_path = await message.download()
        encoded_data_url = encode_image_as_base64(file_path) # Ваша функция кодирования

        # Формируем контент для сохранения в БД (список словарей)
        user_content_list = [{"type": "image_url", "image_url": {"url": encoded_data_url}}]
        caption_text = message.caption or "опиши изображение" # Текст по умолчанию, если нет подписи
        user_content_list.append({"type": "text", "text": caption_text})

        # Вызываем общую функцию обработки
        await process_message(message, user_content_list)

    except Exception as e:
        await message.reply_text(f"❌ Ошибка при обработке изображения: {e}")
    finally:
        if file_path and os.path.exists(file_path):
            os.remove(file_path)

@client.on_message(filters.command("gen"))
async def ask_prompt(_, message: Message):
    user_states[message.chat.id] = "awaiting_prompt"

    cancel_button = InlineKeyboardMarkup([
        [InlineKeyboardButton("❌ Отменить", callback_data="cancel_gen")]
    ])

    await message.reply_text("🎨 Введи промпт для генерации изображения.", reply_markup=cancel_button)

@client.on_callback_query(filters.regex("cancel_context"))
async def cancel_context(_, query):
    chat_id = query.message.chat.id
    user_states.pop(chat_id, None)
    await query.message.delete()

@client.on_callback_query(filters.regex("cancel_gen"))
async def cancel_gen(_, query):
    chat_id = query.message.chat.id
    user_states.pop(chat_id, None)
    await query.message.delete()



# Команда /info
@client.on_message(filters.command("info"))
async def info(_, message: Message):
    chat_id = str(message.chat.id)
    db = next(get_db())
    try:
        chat = db.execute(select(Chat).filter(Chat.chat_id == chat_id)).scalar_one_or_none()
        if chat is None:
            await message.reply_text("❗ Чат не зарегистрирован. Напиши /start.")
            return

        model = chat.model_name
        system_prompt = chat.system_prompt or "⚠️ Не установлен." # Используем значение из БД

        text = (
            f"📊 <b>Информация о текущем чате</b>\n\n"
            f"🤖 <b>Модель:</b> <code>{model}</code>\n"
            f"🧠 <b>Системный промпт:</b>\n<code>{system_prompt}</code>"
        )
        await message.reply_text(text, parse_mode=ParseMode.HTML)
    finally:
        db.close()

@client.on_message(filters.command("help"))
async def help_command(_, message: Message):
    help_text = (
        "🤖 <b>GPT Telegram Бот — что я умею</b>\n\n"

        "📌 <b>/start</b> — Регистрация чата\n"
        "➤ Обязательный шаг перед использованием бота\n"
        "📝 Пример: <code>/start</code>\n\n"

        "🧠 <b>/context</b> — Установить системный стиль общения\n"
        "➤ Задаёт боту поведение (например, «будь экспертом по психологии»)\n"
        "📝 Пример: <code>/context Ты — весёлый бард, говоришь стихами</code>\n\n"

        "🧹 <b>/forget</b> — Очистить историю общения\n"
        "➤ Удаляет контекст предыдущих сообщений\n"
        "📝 Пример: <code>/forget</code>\n\n"

        "🤖 <b>/model</b> — Выбор модели GPT\n"
        "➤ Нажми на кнопку с нужной моделью\n"
        "📝 Пример: <code>/model</code>\n\n"

        "🎨 <b>/gen</b> — Генерация изображения по описанию\n"
        "➤ Введи команду, затем отправь описание картинки\n"
        "📝 Пример:\n<code>/gen</code>\n<code>кошка на подоконнике</code>\n\n"

        "ℹ️ <b>/info</b> — Показать текущую модель и системный промпт\n"
        "📝 Пример: <code>/info</code>\n\n"

        "🆘 <b>/help</b> — Это справочное меню\n"
        "📝 Пример: <code>/help</code>\n\n"

        "<i>Дополнительно:</i>\n"
        "📂 Можешь отправлять файлы .txt, .docx, .pptx, .fb2 — бот прочитает и ответит\n"
        "🖼 Фото тоже можно — бот опишет изображение и даст ответ\n"
    )
    await message.reply_text(help_text, parse_mode=ParseMode.HTML)



# Запуск бота
print("🤖 GPT Telegram бот запущен...")
client.run()